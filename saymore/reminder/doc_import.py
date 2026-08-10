"""文件/图片导入——从给定路径抽出纯文字，交给提醒对话（选文件框在主界面「导入」tab）。

抽字纯代码：docx/xlsx/pptx 用对应库，txt 直读，图片走 RapidOCR（离线、无显存）。
各类型的库都延迟导入：没装/没用到不影响其他类型。
"""
from pathlib import Path

_OCR = None  # RapidOCR 实例惰性加载（首次用才建，避免拖慢启动）


def extract(path):
    """按扩展名抽出纯文字。不支持的类型抛 ValueError。"""
    ext = Path(path).suffix.lower()
    if ext == ".txt":
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    if ext == ".docx":
        return _docx(path)
    if ext == ".xlsx":
        return _xlsx(path)
    if ext == ".pptx":
        return _pptx(path)
    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".webp"):
        return _ocr(path)
    raise ValueError(f"不支持的文件类型: {ext}")


def _docx(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:  # 表格里常放日程，逐格取
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _xlsx(path):
    import openpyxl
    from openpyxl.utils.datetime import from_excel
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    # ponytail: 日历常把日期存成"General"格式的序列号（46167=2026-05-25），
    # openpyxl 只能原样给整数。落在合理日期区间（约 2000–2099）的整数补转回日期，
    # 否则 LLM 拿到一堆 46167 没法建提醒；非日期返回 None。
    def asdate(v):
        if isinstance(v, int) and 36526 <= v <= 73050:
            return from_excel(v, wb.epoch).date().isoformat()
        return None

    # 网格课表里"哪节课在哪天"靠的是列对齐：某行是一排日期（一周表头），下面几行
    # 把课摆在对应星期的列下。直接 [v for v in row if v]会塌缩列、丢掉日期归属，
    # 导致 LLM 只能瞎猜日期。故记住每列当前所属日期，把课粘回它那列的日期。
    lines = []
    for ws in wb.worksheets:
        col_dates = {}  # 列号 -> 该列当前日期
        for row in ws.iter_rows(values_only=True):
            dates = {i: asdate(v) for i, v in enumerate(row) if asdate(v)}
            if len(dates) >= 2:  # 主要是日期的行 = 日历的"周"表头，刷新各列日期
                col_dates = dates
                lines.append(" | ".join(dates[i] for i in sorted(dates)))
                continue
            cells = []
            for i, v in enumerate(row):
                if v is None:
                    continue
                d = col_dates.get(i)
                cells.append(f"{d} {v}" if d else str(v))
            if cells:
                lines.append(" | ".join(cells))
    wb.close()
    return "\n".join(lines)


def _pptx(path):
    import pptx
    prs = pptx.Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _ocr(path):
    global _OCR
    if _OCR is None:
        from rapidocr_onnxruntime import RapidOCR
        _OCR = RapidOCR()
    result, _ = _OCR(path)
    if not result:
        return ""
    return "\n".join(line[1] for line in result)  # 每项是 [box, text, score]


def _selfcheck():
    import tempfile, os
    p = os.path.join(tempfile.gettempdir(), "doc_import_selftest.txt")
    Path(p).write_text("提醒我明天下午三点开会\n后天交报告", encoding="utf-8")
    assert "开会" in extract(p)
    os.remove(p)
    try:
        extract("x.unknown")
        assert False, "应对未知类型报错"
    except ValueError:
        pass
    print("doc_import 自检通过（txt + 未知类型；docx/xlsx/pptx/图片需真实文件手测）。")


if __name__ == "__main__":
    _selfcheck()
